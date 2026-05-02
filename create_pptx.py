"""
Build the OpenLineage Confluent deck on top of the Confluent brand template.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = "/Users/vahidfereydounkolahi/openlineage-confluent/openlineage-confluent-deck.pptx"

# ── Brand palette ─────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x17, 0x33, 0x61)   # #173361  Confluent dark navy
DKNAVY  = RGBColor(0x0D, 0x18, 0x2B)   # #0D182B  darkest navy
BLUE    = RGBColor(0x00, 0xA4, 0xE4)   # #00A4E4  Confluent blue
TEAL    = RGBColor(0x01, 0xCE, 0xDB)   # #01CEDB  Confluent teal
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT  = RGBColor(0xF6, 0xF7, 0xF7)  # #F6F7F7  near-white
LTGRAY  = RGBColor(0xF3, 0xF3, 0xF3)  # #F3F3F3  slide bg
MDGRAY  = RGBColor(0x82, 0x94, 0x94)  # muted
TEXT    = RGBColor(0x17, 0x33, 0x61)  # body text
GREEN   = RGBColor(0x11, 0x50, 0x00)  # #115000  dark green accent
LTGREEN = RGBColor(0xE6, 0xF5, 0xFB)  # #E6F5FB  very light blue
LTBLUE  = RGBColor(0xE0, 0xE6, 0xFF)  # light lavender
ORANGE  = RGBColor(0xAE, 0x1B, 0x85)  # Confluent magenta/accent

# Slide dimensions (10" × 5.62")
SW = Inches(10.0)
SH = Inches(5.62)
FONT = "Inter"


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # blank layout
    return prs.slides.add_slide(layout)


def add_rect(slide, l, t, w, h, fill: RGBColor, *, line_color=None, line_pt=0.5):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h, *,
                size=12, bold=False, color=TEXT,
                align=PP_ALIGN.LEFT, font=FONT,
                wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_multiline(slide, lines, l, t, w, h, *, wrap=True):
    """lines: list of (text, size, bold, color, align)"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    first = True
    for (text, size, bold, color, align) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after  = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.name  = FONT
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
    return txb


# ── Composite components ──────────────────────────────────────────────────────

def slide_header(slide, title, subtitle=""):
    """Navy header bar with title + optional subtitle."""
    add_rect(slide, 0, 0, SW, Inches(1.72), NAVY)
    add_rect(slide, 0, 0, SW, Pt(5), TEAL)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.18), Inches(9.2), Inches(0.75),
                size=22, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.65),
                    size=11, color=OFFWHT)


def card(slide, title, body, l, t, w, h, accent=BLUE, *, body_size=10):
    add_rect(slide, l, t, w, h, WHITE, line_color=RGBColor(0xD8, 0xD8, 0xD8))
    add_rect(slide, l, t, Pt(4), h, accent)
    add_multiline(slide, [
        (title, 11, True, NAVY, PP_ALIGN.LEFT),
        ("", 3,  False, TEXT, PP_ALIGN.LEFT),
        (body, body_size, False, TEXT, PP_ALIGN.LEFT),
    ], l + Pt(9), t + Pt(5), w - Pt(14), h - Pt(8))


def stat_box(slide, num, lbl, l, t, w):
    h = Inches(1.0)
    add_rect(slide, l, t, w, h, WHITE, line_color=RGBColor(0xD8, 0xD8, 0xD8))
    add_textbox(slide, num,
                l, t + Inches(0.04), w, Inches(0.55),
                size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, lbl,
                l, t + Inches(0.6), w, Inches(0.35),
                size=8, bold=True, color=MDGRAY, align=PP_ALIGN.CENTER)


def simple_table(slide, headers, rows, col_widths_in, top, left=Inches(0.4)):
    rh = Inches(0.4)
    hh = Inches(0.33)
    total_w = Inches(sum(col_widths_in))

    add_rect(slide, left, top, total_w, hh, NAVY)
    x = left
    for hdr, cw in zip(headers, col_widths_in):
        add_textbox(slide, hdr, x + Pt(4), top + Pt(3),
                    Inches(cw) - Pt(6), hh,
                    size=8, bold=True, color=WHITE)
        x += Inches(cw)

    for ri, row in enumerate(rows):
        ry  = top + hh + ri * rh
        bg  = LTGRAY if ri % 2 == 0 else WHITE
        add_rect(slide, left, ry, total_w, rh, bg)
        x = left
        for ci, (cell, cw) in enumerate(zip(row, col_widths_in)):
            add_textbox(slide, cell,
                        x + Pt(4), ry + Pt(3),
                        Inches(cw) - Pt(8), rh - Pt(4),
                        size=9, bold=(ci == 0),
                        color=NAVY if ci == 0 else TEXT)
            x += Inches(cw)


# ── Slide builders ────────────────────────────────────────────────────────────

def s01_title(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, DKNAVY)
    add_rect(s, 0, 0, SW, Pt(5), TEAL)
    add_textbox(s, "CONFLUENT  ·  OPENLINEAGE",
                Inches(0.5), Inches(0.8), Inches(9), Inches(0.4),
                size=10, bold=True, color=TEAL)
    add_multiline(s, [
        ("End-to-End Data Lineage", 36, True, WHITE, PP_ALIGN.LEFT),
        ("for Confluent Cloud",     36, True, TEAL,  PP_ALIGN.LEFT),
    ], Inches(0.5), Inches(1.3), Inches(9), Inches(2.0))
    add_textbox(s,
        "A bridge that translates every Confluent component — managed connectors, "
        "Flink statements, ksqlDB queries, consumer groups, and self-managed Connect "
        "— into the open OpenLineage standard.",
        Inches(0.5), Inches(3.4), Inches(8.5), Inches(1.1),
        size=13, color=OFFWHT)
    add_textbox(s,
        "Managed Connectors  ·  Apache Flink  ·  ksqlDB  ·  Consumer Groups  "
        "·  Self-Managed Connect  ·  Schema Registry",
        Inches(0.5), Inches(4.6), Inches(9.0), Inches(0.5),
        size=10, color=TEAL)


def s02_challenge(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "The Data Lineage Challenge",
                 "Modern data platforms are built on Kafka — but lineage stops at the topic boundary")
    cw = Inches(4.7); ch = Inches(1.7); t1 = Inches(1.85); t2 = t1 + ch + Inches(0.12)
    card(s, "Lineage Blind Spots",
         "Managed connectors, Flink jobs, and ksqlDB queries each live in separate control planes with no shared lineage model.",
         Inches(0.4), t1, cw, ch, ORANGE)
    card(s, "Disconnected Islands",
         "A message flowing from a DatagenSource → Flink → HTTP Sink spans three systems with zero automatic end-to-end tracing.",
         Inches(0.4) + cw + Inches(0.2), t1, cw, ch, ORANGE)
    card(s, "No Standard Format",
         "Each Confluent API returns a different shape. There is no single output format consumable by downstream governance tools.",
         Inches(0.4), t2, cw, ch, BLUE)
    card(s, "Hidden Producer Identity",
         "Kafka producers (Java clients, microservices) leave no fingerprint at the broker level — making self-service lineage impossible.",
         Inches(0.4) + cw + Inches(0.2), t2, cw, ch, BLUE)


def s03_openlineage(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "OpenLineage: The Universal Standard",
                 "A vendor-neutral, open specification for capturing data lineage events at runtime")
    cw3 = Inches(3.0); ch3 = Inches(1.4); t3 = Inches(1.85)
    card(s, "Event-Driven Model",
         "Every job run emits START / COMPLETE / FAIL events. Lineage is captured in real time, not reconstructed from logs.",
         Inches(0.4), t3, cw3, ch3, TEAL)
    card(s, "Jobs + Datasets",
         "Jobs (connectors, Flink statements) consume InputDatasets and produce OutputDatasets. Edges are implicit from events.",
         Inches(0.4) + cw3 + Inches(0.15), t3, cw3, ch3, TEAL)
    card(s, "Extensible via Facets",
         "Core spec is lean. Any system can add custom facets (schema, throughput, SLA) without breaking consumers.",
         Inches(0.4) + (cw3 + Inches(0.15)) * 2, t3, cw3, ch3, NAVY)

    # Flow diagram
    fnodes = [
        ("Job Run\nFlink / Connect / ksqlDB", BLUE),
        ("RunEvent\nSTART / COMPLETE / FAIL",  RGBColor(0x14, 0x44, 0x9A)),
        ("Lineage Backend\nMarquez / DataHub",  NAVY),
        ("Governance UI\nSearch · Impact · Audit", DKNAVY),
    ]
    nw = Inches(2.1); nh = Inches(0.9); fy = Inches(3.55)
    for i, (lbl, col) in enumerate(fnodes):
        fx = Inches(0.4) + i * (nw + Inches(0.3))
        add_rect(s, fx, fy, nw, nh, col)
        add_textbox(s, lbl, fx + Pt(5), fy + Pt(6), nw - Pt(8), nh - Pt(10),
                    size=10, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(s, "→", fx + nw + Pt(3), fy + Pt(16), Inches(0.28), Inches(0.4),
                        size=16, color=MDGRAY, align=PP_ALIGN.CENTER)
    add_textbox(s,
        "Any job type that reads/writes Kafka topics can emit OpenLineage events — "
        "bridging all components into a single unified lineage graph.",
        Inches(0.4), Inches(4.65), Inches(9.2), Inches(0.5),
        size=10, color=MDGRAY, italic=True)


def s04_architecture(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "Bridge Architecture",
                 "A lightweight Python service that polls Confluent APIs and translates to OpenLineage events")
    anodes = [
        ("Confluent APIs\nConnect · Flink\nMetrics · ksqlDB\nKafka REST · SR", DKNAVY),
        ("LineageGraph\nUnified edge model", BLUE),
        ("OL Mapper\nNamespace · Facets\nRunIds", RGBColor(0x14, 0x44, 0x9A)),
        ("Emitter\nDiff-tracking\nParallel · ABORT", NAVY),
        ("Marquez\nor any OL backend", RGBColor(0x04, 0x04, 0x53)),
    ]
    anw = Inches(1.72); agh = Inches(0.13); ay = Inches(1.9); ah = Inches(1.3)
    for i, (lbl, col) in enumerate(anodes):
        ax = Inches(0.4) + i * (anw + agh + Inches(0.14))
        add_rect(s, ax, ay, anw, ah, col)
        add_textbox(s, lbl, ax + Pt(5), ay + Pt(5), anw - Pt(8), ah - Pt(8),
                    size=9, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(s, "→", ax + anw + Pt(3), ay + Pt(22),
                        Inches(0.12), Inches(0.4),
                        size=14, color=MDGRAY, align=PP_ALIGN.CENTER)

    cw4 = Inches(3.0); ch4 = Inches(1.45); t4 = Inches(3.38)
    card(s, "Diff-Tracking State",
         "SQLite state store detects added and removed jobs. Removed jobs emit ABORT events — no stale lineage.",
         Inches(0.4), t4, cw4, ch4, BLUE)
    card(s, "Parallel Emission",
         "All RunEvents emitted concurrently via a thread pool — scales to 10 000+ topics without blocking.",
         Inches(0.4) + cw4 + Inches(0.2), t4, cw4, ch4, TEAL)
    card(s, "Scheduled Pipeline",
         "APScheduler drives periodic polling. run-once for snapshots, run for continuous operation.",
         Inches(0.4) + (cw4 + Inches(0.2)) * 2, t4, cw4, ch4, NAVY)


def s05_sources(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "Five Lineage Sources",
                 "Complete coverage of every Confluent component that moves data")
    simple_table(s,
        ["Source", "API Used", "What It Captures", "Job Type"],
        [
            ["Managed Connect",      "api.confluent.cloud/connect/v1",          "All managed source & sink connectors in Confluent Cloud",            "SOURCE / SINK_CONNECTOR"],
            ["Flink SQL",            "Flink REST API + SQL parser",             "Persistent statements; input/output topics parsed from SQL",          "QUERY"],
            ["Consumer Groups",      "Metrics API consumer_lag_offsets",        "Any client committing offsets — the only cloud-plane signal",          "CONSUMER_GROUP"],
            ["ksqlDB",               "ksqlDB REST POST /ksql",                  "Persistent CREATE STREAM / TABLE queries and their topic wiring",      "QUERY"],
            ["Self-Managed Connect", "GET {endpoint}/connectors?expand=info,…", "On-prem / self-hosted Connect clusters alongside cloud topology",      "SOURCE / SINK_CONNECTOR"],
        ],
        [1.9, 2.45, 3.5, 1.75],
        top=Inches(1.85),
    )
    sw5 = Inches(2.15); sy = Inches(4.42)
    stat_box(s, "5",    "Lineage Sources",        Inches(0.4),              sy, sw5)
    stat_box(s, "146",  "Tests Passing",          Inches(0.4) + sw5 + Inches(0.17),      sy, sw5)
    stat_box(s, "0",    "Credentials for Tests",  Inches(0.4) + (sw5+Inches(0.17))*2,   sy, sw5)
    stat_box(s, "10k+", "Topics Supported",       Inches(0.4) + (sw5+Inches(0.17))*3,   sy, sw5)


def s06_topology(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "Live Demo Topology",
                 "End-to-end lineage stitching four job types into a single directed graph")

    CTEAL = RGBColor(0xFF, 0xF3, 0xE0)  # connect (orange tint)
    CBLU  = RGBColor(0xE6, 0xF5, 0xFB)  # topic (blue tint)
    CGRN  = RGBColor(0xE6, 0xF5, 0xE9)  # flink (green tint)
    CRED  = RGBColor(0xFB, 0xE9, 0xE7)  # java producer
    CPUR  = RGBColor(0xF3, 0xE5, 0xF5)  # consumer group

    def tn(label, bg, l, t, w=Inches(2.35), h=Inches(0.44)):
        add_rect(s, l, t, w, h, bg, line_color=RGBColor(0xD8, 0xD8, 0xD8))
        add_textbox(s, label, l + Pt(5), t + Pt(4), w - Pt(8), h - Pt(6),
                    size=8, bold=True, color=TEXT)

    def ar(sym, l, t):
        add_textbox(s, sym, l, t, Inches(0.3), Inches(0.35),
                    size=12, color=MDGRAY, align=PP_ALIGN.CENTER)

    tn("Datagen Source Connector", CTEAL, Inches(0.3),  Inches(1.85))
    tn("ol-raw-orders",            CBLU,  Inches(2.9),  Inches(1.85), w=Inches(1.9))
    tn("Java OrderProducer",       CRED,  Inches(5.05), Inches(1.85))
    ar("→", Inches(2.68),  Inches(1.95))
    ar("←", Inches(4.76),  Inches(1.95))

    tn("Flink: ol-orders-enricher", CGRN, Inches(1.0), Inches(2.52))
    tn("ol-orders-enriched",        CBLU, Inches(3.6), Inches(2.52), w=Inches(2.0))
    ar("↓", Inches(3.75),  Inches(2.35))
    ar("→", Inches(3.37),  Inches(2.64))

    tn("Flink: ol-high-value-alerts",   CGRN, Inches(0.3),  Inches(3.2))
    tn("ol-high-value-alerts",          CBLU, Inches(2.9),  Inches(3.2), w=Inches(1.8))
    tn("Flink: ol-medium-risk-orders",  CGRN, Inches(5.05), Inches(3.2))
    tn("ol-medium-risk-orders",         CBLU, Inches(7.65), Inches(3.2), w=Inches(2.0))
    ar("↓", Inches(4.52),  Inches(2.99))
    ar("→", Inches(2.68),  Inches(3.32))
    ar("→", Inches(7.42),  Inches(3.32))

    tn("HTTP Sink Connector",           CTEAL, Inches(1.0),  Inches(3.88))
    tn("Consumer Group: connect-lcc-*", CPUR,  Inches(4.3),  Inches(3.88), w=Inches(2.7))
    ar("←", Inches(3.37),  Inches(4.0))
    ar("←", Inches(4.05),  Inches(4.0))
    ar("↑", Inches(4.85),  Inches(3.68))

    add_textbox(s,
        "Orange = Connect    Blue = Kafka Topic    Green = Flink    Purple = Consumer Group    Red = Java Producer",
        Inches(0.3), Inches(4.64), Inches(9.4), Inches(0.35),
        size=8, color=MDGRAY, italic=True)


def s07_metadata(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "Rich Topic Metadata via Custom Facets",
                 "Topics carry schema, physical config, and live throughput — all queryable in the lineage backend")

    lw = Inches(4.2); ch7 = Inches(1.18); t7 = Inches(1.85)
    card(s, "Schema Facet (SchemaDatasetFacet)",
         "Key and value subjects from Schema Registry merged into key.* / value.* prefixed fields. Supports AVRO, JSON Schema, Protobuf.",
         Inches(0.4), t7, lw, ch7, BLUE)
    card(s, "KafkaTopicDatasetFacet",
         "Partition count, replication factor, and internal-topic flag fetched from Kafka REST. Mirrors Confluent Stream Lineage topic nodes.",
         Inches(0.4), t7 + ch7 + Inches(0.12), lw, ch7, TEAL)
    card(s, "KafkaTopicThroughputDatasetFacet",
         "Bytes in/out and records in/out over a configurable lookback window from the Metrics API. Enables data-volume impact analysis.",
         Inches(0.4), t7 + (ch7 + Inches(0.12)) * 2, lw, ch7, NAVY)

    code = (
        "# Facets on ol-orders-enriched\n\n"
        "schema:\n"
        "  fields:\n"
        "    - name: value.order_id   type: string\n"
        "    - name: value.amount     type: double\n"
        "    - name: value.customer   type: string\n"
        "    - name: key.order_id     type: string\n\n"
        "kafkaTopic:\n"
        "  partitions:        6\n"
        "  replicationFactor: 3\n\n"
        "kafkaThroughput:\n"
        "  bytesIn:       1_482_304\n"
        "  recordsIn:     12_400\n"
        "  windowMinutes: 10"
    )
    rw = SW - lw - Inches(0.85)
    add_rect(s, lw + Inches(0.65), Inches(1.85), rw, Inches(3.65), DKNAVY)
    add_textbox(s, code,
                lw + Inches(0.78), Inches(1.98), rw - Inches(0.22), Inches(3.45),
                size=9, color=OFFWHT, font="Courier New")


def s08_namespaces(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "Namespace Conventions",
                 "Every job and dataset receives a stable, globally unique identifier in the OpenLineage graph")
    simple_table(s,
        ["Component", "OL Job Namespace", "OL Job Name", "processingType"],
        [
            ["Managed Connector",      "kafka-connect://<env_id>",         "<connector_name>",  "STREAMING"],
            ["Self-Managed Connector", "kafka-connect://<cluster_label>",  "<connector_name>",  "STREAMING"],
            ["Flink Statement",        "flink://<env_id>",                 "<statement_name>",  "STREAMING"],
            ["Consumer Group",         "kafka-consumer-group://<cluster>", "<group_id>",        "STREAMING"],
            ["ksqlDB Query",           "ksqldb://<ksql_cluster_id>",       "<query_id>",        "STREAMING"],
            ["Java SDK Producer",      "kafka-producer://<service_name>",  "<job_name>",        "STREAMING"],
            ["Kafka Topic (Dataset)",  "kafka://<bootstrap_server>",       "<topic_name>",      "—"],
        ],
        [2.2, 2.8, 2.2, 1.9],
        top=Inches(1.85),
    )
    add_rect(s, Inches(0.4), Inches(4.48), Inches(9.2), Inches(0.78),
             RGBColor(0xE6, 0xF5, 0xFB), line_color=TEAL)
    add_rect(s, Inches(0.4), Inches(4.48), Pt(4), Inches(0.78), TEAL)
    add_textbox(s,
        "Stable RunIDs are derived deterministically from namespace + name + cycle_key via SHA-256 "
        "— ensuring the same job maps to the same run across polling cycles, preventing duplicate lineage nodes.",
        Inches(0.56), Inches(4.52), Inches(8.9), Inches(0.7),
        size=9, color=NAVY)


def s09_deployment(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "Simple Deployment",
                 "One config file, three CLI commands — no code changes required")

    config = (
        "# config.yaml\n\n"
        "confluent:\n"
        "  api_key:    MRYYEBG4MF6FM2BP\n"
        "  api_secret: ***\n"
        "  env_id:     env-m2qxq\n"
        "  cluster_id: lkc-1j6rd3\n\n"
        "  schema_registry:\n"
        "    endpoint: https://psrc-….confluent.cloud\n"
        "    api_key:  BPYPBFCQOV6A63XQ\n\n"
        "  kafka_rest:\n"
        "    endpoint: https://pkc-….confluent.cloud\n"
        "    api_key:  JBFEOVMJA543R4MK\n\n"
        "openlineage:\n"
        "  backend_url:     http://localhost:5000\n"
        "  kafka_bootstrap: pkc-pgq85…:9092\n"
        "  interval_sec:    60"
    )
    add_rect(s, Inches(0.4), Inches(1.85), Inches(4.65), Inches(3.65), DKNAVY)
    add_textbox(s, config, Inches(0.52), Inches(1.95), Inches(4.42), Inches(3.5),
                size=9, color=OFFWHT, font="Courier New")

    rx = Inches(5.3); rw = Inches(4.35); ch9 = Inches(0.84)
    t9 = Inches(1.85); gp = Inches(0.12)
    card(s, "ol-confluent validate",
         "Checks all credentials and API connectivity before any data is emitted. Fast pre-flight for CI/CD.",
         rx, t9, rw, ch9, BLUE)
    card(s, "ol-confluent run-once",
         "Fetches the full topology snapshot and emits all RunEvents once. Ideal for scheduled jobs or demos.",
         rx, t9 + ch9 + gp, rw, ch9, TEAL)
    card(s, "ol-confluent run",
         "Continuous polling on the configured interval. Detects added and removed jobs automatically.",
         rx, t9 + (ch9 + gp) * 2, rw, ch9, NAVY)
    card(s, "Docker Compose Ready",
         "Marquez backend + Web UI included. One docker compose up -d to start the full stack.",
         rx, t9 + (ch9 + gp) * 3, rw, ch9, MDGRAY)


def s10_results(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "Results at a Glance",
                 "What the bridge delivers out of the box today")

    sw10 = Inches(2.15); sy = Inches(1.85)
    stat_box(s, "5",   "Lineage Sources",       Inches(0.4),               sy, sw10)
    stat_box(s, "4",   "Job Types Identified",  Inches(0.4)+sw10+Inches(0.17),      sy, sw10)
    stat_box(s, "3",   "Custom Facets",         Inches(0.4)+(sw10+Inches(0.17))*2,  sy, sw10)
    stat_box(s, "146", "Tests (all offline)",   Inches(0.4)+(sw10+Inches(0.17))*3,  sy, sw10)

    cw10 = Inches(4.7); ch10 = Inches(1.4); t10 = Inches(3.07)
    card(s, "Full End-to-End Graph",
         "Connector → Topic → Flink → Topic → Sink rendered as a single navigable lineage graph. Depth-4 view shows the complete data flow in one screen.",
         Inches(0.4), t10, cw10, ch10, RGBColor(0x11, 0x50, 0x00))
    card(s, "Schema Propagation",
         "AVRO / JSON Schema / Protobuf field lists from Schema Registry attached to every topic dataset — enabling column-level impact analysis.",
         Inches(0.4) + cw10 + Inches(0.2), t10, cw10, ch10, RGBColor(0x11, 0x50, 0x00))
    card(s, "Stale Lineage Detection",
         "Deleted or stopped jobs emit ABORT events automatically. The lineage graph stays current without manual intervention.",
         Inches(0.4), t10 + ch10 + Inches(0.12), cw10, ch10, BLUE)
    card(s, "Java SDK Integration",
         "The included OrderProducer demo shows how any Java Kafka client can self-instrument START / COMPLETE events using the OpenLineage Java SDK.",
         Inches(0.4) + cw10 + Inches(0.2), t10 + ch10 + Inches(0.12), cw10, ch10, BLUE)


def s11_roadmap(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, LTGRAY)
    slide_header(s, "Roadmap & Next Steps",
                 "Prioritized extensions to expand coverage and ecosystem integration")
    items = [
        ("Flink SQL UDTF Scanner",
         "Extend the SQL parser to extract field-level lineage from Flink UDTFs — enabling column-level impact graphs.", BLUE),
        ("Producer Auto-Discovery",
         "Correlate Kafka client metrics with application metadata to identify individual producers without SDK instrumentation.", BLUE),
        ("DataHub & Apache Atlas",
         "Swap Marquez for DataHub or Apache Atlas by changing a single backend_url line in config.", TEAL),
        ("Helm Chart / Operator",
         "Package as a Kubernetes Helm chart for one-command deployment alongside Confluent Platform or CFK.", TEAL),
        ("Change Notifications",
         "Emit Slack / PagerDuty alerts when a high-value topic loses its producer or consumer lag exceeds a threshold.", NAVY),
        ("Multi-Environment Support",
         "Single bridge instance polling multiple Confluent environments — unified lineage across dev, staging, and prod.", NAVY),
    ]
    cw11 = Inches(3.0); ch11 = Inches(1.7); t11 = Inches(1.85)
    for i, (title, body, acc) in enumerate(items):
        col = i % 3; row = i // 3
        card(s, title, body,
             Inches(0.4) + col * (cw11 + Inches(0.2)),
             t11 + row * (ch11 + Inches(0.12)),
             cw11, ch11, acc)


def s12_thankyou(prs):
    s = new_slide(prs)
    add_rect(s, 0, 0, SW, SH, DKNAVY)
    add_rect(s, 0, 0, SW, Pt(5), TEAL)
    add_textbox(s, "Questions?",
                Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.8),
                size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s,
        "OpenLineage Bridge for Confluent Cloud\n"
        "Complete end-to-end data lineage across every component of your streaming data platform.",
        Inches(1.0), Inches(3.35), Inches(8.0), Inches(1.1),
        size=14, color=OFFWHT, align=PP_ALIGN.CENTER)
    add_textbox(s,
        "openlineage.io   ·   confluent.io/product/flink   ·   marquezproject.ai",
        Inches(1.0), Inches(4.55), Inches(8.0), Inches(0.5),
        size=11, color=TEAL, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    s01_title(prs)
    s02_challenge(prs)
    s03_openlineage(prs)
    s04_architecture(prs)
    s05_sources(prs)
    s06_topology(prs)
    s07_metadata(prs)
    s08_namespaces(prs)
    s09_deployment(prs)
    s10_results(prs)
    s11_roadmap(prs)
    s12_thankyou(prs)

    prs.save(OUTPUT)
    print(f"Saved {len(prs.slides)} slides → {OUTPUT}")


if __name__ == "__main__":
    main()
